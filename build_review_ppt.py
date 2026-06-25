# -*- coding: utf-8 -*-
"""I-Study — 개선사항 & 향후 계획 (PPT 1장)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 색상 (성능평가 보고서 톤과 통일) ──
NAVY   = RGBColor(0x1F, 0x3A, 0x5F)
BLUE   = RGBColor(0x2E, 0x6F, 0xB7)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GREENbg= RGBColor(0xE8, 0xF5, 0xEC)
AMBER  = RGBColor(0xC7, 0x82, 0x12)
AMBERbg= RGBColor(0xFB, 0xF1, 0xDD)
DGRAY  = RGBColor(0x2B, 0x2B, 0x2B)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LGRAY  = RGBColor(0x8A, 0x8A, 0x8A)
STATbg = RGBColor(0xF2, 0xF6, 0xFB)
STATbd = RGBColor(0xDD, 0xE6, 0xF1)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = "Malgun Gothic"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
shp = slide.shapes


def set_font(run, name=FONT, size=None, bold=None, color=None):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)
    if size  is not None: run.font.size = Pt(size)
    if bold  is not None: run.font.bold = bold
    if color is not None: run.font.color.rgb = color


def rect(shape_type, x, y, w, h, fill=None, line=None, lw=1.0, radius=None):
    s = shp.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    if radius is not None:
        try: s.adjustments[0] = radius
        except Exception: pass
    return s


def textbox(x, y, w, h, valign=MSO_ANCHOR.TOP):
    tb = shp.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return tf


# ════════════ 제목 영역 ════════════
pill = rect(MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, 0.42, 3.05, 0.42,
            fill=NAVY, radius=0.5)
ptf = pill.text_frame; ptf.word_wrap = False
ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
pr = ptf.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
set_font(pr.add_run(), size=11.5, bold=True, color=WHITE); pr.runs[0].text = "I-STUDY  ·  프로젝트 리뷰"

tf = textbox(0.52, 0.95, 9.5, 0.72)
r = tf.paragraphs[0].add_run(); r.text = "개선사항 & 향후 계획"
set_font(r, size=34, bold=True, color=NAVY)

tf = textbox(0.55, 1.66, 9.0, 0.4)
r = tf.paragraphs[0].add_run(); r.text = "AI 시선추적 학습 집중도 관리 시스템"
set_font(r, size=14, color=GRAY)

tf = textbox(9.7, 0.5, 3.1, 0.7, valign=MSO_ANCHOR.MIDDLE)
for i, t in enumerate(["한이음 프로젝트 · 2026", "데스크톱 앱 + 웹 대시보드"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.RIGHT
    rr = p.add_run(); rr.text = t
    set_font(rr, size=10.5, color=LGRAY)

# ════════════ 컬럼 설정 ════════════
COL_TOP = 2.18
COL_W   = 5.96
LEFT_X  = 0.55
RIGHT_X = 6.82
HDR_H   = 0.5
ITEM_TOP = COL_TOP + HDR_H + 0.2
ITEM_H  = 0.605
ITEM_STEP = 0.625

improvements = [
    ("백엔드 아키텍처 분리", "REST·웹은 Spring Boot, AI 시선추적은 FastAPI로 이원화"),
    ("데이터 모델 통합",     "study_sessions 단일화로 앱·웹 학습기록 통합 (중복 4개 제거)"),
    ("데스크톱 앱 경량화",   "집중 모드에 집중 — 통계·화이트리스트는 웹으로 이전"),
    ("AI 인식률 검증",       "AFLW2000 4,000장으로 시선 방향 인식 97.6% 달성"),
    ("손쉬운 배포",          "설치 불필요 EXE + LAN·오프라인 모드 지원"),
]
future = [
    ("AI 정확도 고도화",     "취약한 '아래' 방향(84%) 데이터 보강·ML 모델 상시 적용"),
    ("캘리브레이션 자동화",  "사용자별 시선 임계값 자동 보정으로 정확도 향상"),
    ("클라우드 배포",        "LAN 한정 → 어디서나 접속 가능한 클라우드 호스팅"),
    ("학습 리포트 강화",     "집중도 추이 분석 및 지도자용 자동 리포트 생성"),
    ("알림·플랫폼 확장",     "주간 리포트·푸시 알림, 모바일 대응 검토"),
]


def column(x, header, items, color, color_bg, glyph):
    # 헤더 카드
    rect(MSO_SHAPE.ROUNDED_RECTANGLE, x, COL_TOP, COL_W, HDR_H, fill=color, radius=0.28)
    htf = textbox(x + 0.28, COL_TOP, COL_W - 0.4, HDR_H, valign=MSO_ANCHOR.MIDDLE)
    hp = htf.paragraphs[0]
    rg = hp.add_run(); rg.text = glyph + "  "; set_font(rg, size=16, bold=True, color=WHITE)
    rh = hp.add_run(); rh.text = header;        set_font(rh, size=16, bold=True, color=WHITE)
    # 항목
    for i, (kw, desc) in enumerate(items):
        iy = ITEM_TOP + i * ITEM_STEP
        # 번호 원
        d = 0.34
        circ = rect(MSO_SHAPE.OVAL, x + 0.04, iy + (ITEM_H - d) / 2, d, d, fill=color_bg)
        ctf = circ.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = str(i + 1); set_font(cr, size=12, bold=True, color=color)
        # 텍스트
        t = textbox(x + 0.52, iy, COL_W - 0.56, ITEM_H, valign=MSO_ANCHOR.MIDDLE)
        p = t.paragraphs[0]; p.space_after = Pt(0); p.line_spacing = 1.04
        r1 = p.add_run(); r1.text = kw;          set_font(r1, size=12.5, bold=True, color=color)
        r2 = p.add_run(); r2.text = "  " + desc; set_font(r2, size=11, color=GRAY)


column(LEFT_X,  "개선사항 (완료)", improvements, GREEN, GREENbg, "✓")
column(RIGHT_X, "향후 계획",       future,       AMBER, AMBERbg, "→")

# ════════════ 하단 핵심 지표 ════════════
STAT_TOP = 6.22
STAT_H   = 0.84
stats = [
    ("97.6%", "AI 시선방향 인식 정확도", GREEN),
    ("30 FPS", "실시간 시선추적 처리 성능", BLUE),
    ("7개", "통합 DB 테이블 (중복 4개 제거)", NAVY),
]
sx = 0.55; sgap = 0.3
sw = (12.78 - 0.55 - 2 * sgap) / 3
for i, (num, label, col) in enumerate(stats):
    cx = sx + i * (sw + sgap)
    rect(MSO_SHAPE.ROUNDED_RECTANGLE, cx, STAT_TOP, sw, STAT_H,
         fill=STATbg, line=STATbd, lw=1.0, radius=0.16)
    nt = textbox(cx + 0.26, STAT_TOP + 0.1, sw - 0.5, 0.42, valign=MSO_ANCHOR.MIDDLE)
    nr = nt.paragraphs[0].add_run(); nr.text = num; set_font(nr, size=24, bold=True, color=col)
    lt = textbox(cx + 0.28, STAT_TOP + 0.5, sw - 0.5, 0.3, valign=MSO_ANCHOR.MIDDLE)
    lr = lt.paragraphs[0].add_run(); lr.text = label; set_font(lr, size=10.5, color=GRAY)

OUT = "I-Study_개선사항_향후계획.pptx"
prs.save(OUT)
print("saved:", OUT)
